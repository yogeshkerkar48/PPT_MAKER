"""
PowerPoint Generator Module - Cinematic V4.2
Adaptive Aesthetics: Theme-aware cards and ultra-tight sizing.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from typing import Tuple, Dict
from io import BytesIO
import os
from PIL import Image

def hex_to_rgb(hex_color: str, default_color: Tuple[int, int, int] = (25, 42, 86)) -> Tuple[int, int, int]:
    if not hex_color:
        return default_color
    hex_color = hex_color.lstrip('#')
    if len(hex_color) != 6:
        return default_color
    try:
        return (int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    except:
        return default_color

def setup_slide_layout(prs, slide, bg_color, image_path, index=0):
    """
    Sets up the visual layout with aspect-ratio aware image scaling.
    Alternates left/right based on index.
    """
    # 1. Solid Background
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
    
    # Updated Config for Larger, Consistent Presence (Matching User Screenshot)
    IMAGE_MAX_WIDTH = prs.slide_width * 0.55
    IMAGE_MAX_HEIGHT = prs.slide_height * 0.85
    TEXT_BOX_WIDTH = prs.slide_width * 0.4
    
    # Alternating Logic
    is_left_style = (index % 2 == 0) # Even INDEX: Image Left, Text Right
    
    if is_left_style:
        # Text on Right
        text_area = {
            'left': prs.slide_width - TEXT_BOX_WIDTH,
            'top': Inches(0.5),
            'width': TEXT_BOX_WIDTH - Inches(0.2)
        }
        img_left_pos = Inches(0.5)
    else:
        # Text on Left
        text_area = {
            'left': Inches(0.2),
            'top': Inches(0.5),
            'width': TEXT_BOX_WIDTH - Inches(0.2)
        }
        img_left_pos = TEXT_BOX_WIDTH + Inches(0.2)

    if image_path and os.path.exists(image_path) and os.path.getsize(image_path) > 0:
        try:
            # Get original dimensions
            with Image.open(image_path) as img:
                orig_width, orig_height = img.size
            aspect_ratio = orig_width / orig_height
            
            # Calculate scaled dimensions (preserve aspect ratio)
            if aspect_ratio > 1:  # Wide image
                scaled_width = min(IMAGE_MAX_WIDTH, prs.slide_width * 0.9)
                scaled_height = scaled_width / aspect_ratio
                if scaled_height > IMAGE_MAX_HEIGHT:
                    scaled_height = IMAGE_MAX_HEIGHT
                    scaled_width = scaled_height * aspect_ratio
            else:  # Tall image
                scaled_height = min(IMAGE_MAX_HEIGHT, prs.slide_height * 0.8)
                scaled_width = scaled_height * aspect_ratio
            
            # Position adjustment for right-aligned image
            final_img_left = img_left_pos
            if not is_left_style:
                # If image is on right, we can center it in the remaining 60%
                available_width = prs.slide_width - TEXT_BOX_WIDTH
                final_img_left = TEXT_BOX_WIDTH + (available_width - scaled_width) / 2

            # Add image to slide
            slide.shapes.add_picture(
                image_path,
                left=final_img_left,
                top=(prs.slide_height - scaled_height) / 2, # Vertically centered
                width=scaled_width,
                height=scaled_height
            )
        except Exception as e:
            print(f"Image Placement Error: {e}")
            
    return text_area

def add_adaptive_card(prs, slide, title_text, bullet_text, top, width, left_pos, accent_hex='6366f1'):
    """
    Adds a text card at specific coordinates.
    """
    # 1. Add the Shape
    txBox = slide.shapes.add_textbox(left_pos, top, width, prs.slide_height * 0.8)
    
    # 2. Enable Auto-Size
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # 3. STRONG Semi-transparent dark background
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = RGBColor(0, 0, 0)
    txBox.fill.transparency = 0.2
    txBox.line.fill.background()
    
    # --- Adaptive Colors (Restore Cinematic Look) ---
    accent_rgb = hex_to_rgb(accent_hex)
    # Brightness heuristic for title safety
    accent_brightness = (accent_rgb[0] * 299 + accent_rgb[1] * 587 + accent_rgb[2] * 114) / 1000
    title_color = RGBColor(*accent_rgb) if accent_brightness > 100 else RGBColor(255, 255, 0)
    
    # 5. Bring to Front
    try:
        slide.shapes._spTree.remove(txBox._element)
        slide.shapes._spTree.append(txBox._element)
    except:
        pass
    
    # 6. Content Margins
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    
    # --- TITLE ---
    p_title = tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.size = Pt(28) # Restored to larger size
    p_title.font.bold = True
    p_title.font.color.rgb = title_color # Restored colorful title
    p_title.space_after = Pt(12)
        
    # --- BULLET POINTS ---
    if bullet_text:
        for line in bullet_text.split('\n'):
            if not line.strip(): continue
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18) # Restored to larger size
            p.font.color.rgb = RGBColor(255, 255, 255) # Pure white
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(8)
            p.level = 0
        
    return txBox

def create_slide(prs, slide_data, bg_color, image_path, is_title=False, index=0):
    """Creates a slide with alternating left/right positioning."""
    try:
        slide_layout = prs.slide_layouts[6] # Blank
    except:
        slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    accent = slide_data.get('accent_color', '6366f1')

    # Setup Layout & Image (Returns valid text area)
    text_area = setup_slide_layout(prs, slide, bg_color, image_path, index=index)
    
    # Prepare Content
    if is_title:
        bullets = slide_data.get('bullet_points', [])
        bullet_text = "\n".join([str(p) for p in bullets])
        # Title slide text at top
        add_adaptive_card(prs, slide, slide_data['title'], bullet_text, Inches(0.5), text_area['width'], text_area['left'], accent_hex=accent)
    else:
        bullets = slide_data.get('bullet_points', [])
        bullet_text = "\n".join([f"• {p}" for p in bullets])
        add_adaptive_card(prs, slide, slide_data['title'], bullet_text, Inches(1.5), text_area['width'], text_area['left'], accent_hex=accent)

def generate_pptx(slides_data: Dict, bg_color: str, slide_images: list = None) -> BytesIO:
    prs = Presentation()
    rgb_color = hex_to_rgb(bg_color)
    slides = slides_data.get('slides', [])
    
    for i, slide_data in enumerate(slides):
        img_path = slide_images[i] if slide_images and i < len(slide_images) else None
        # Use index to alternate layout
        create_slide(prs, slide_data, rgb_color, img_path, is_title=False, index=i)
    
    pptx_buffer = BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    return pptx_buffer
